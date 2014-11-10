from BaseHTTPServer import BaseHTTPRequestHandler, HTTPServer
import urlparse

def generate_ppt(name):
    import pptx
    from pptx.util import Inches, Pt
    f = open("blue.pptx")
    prs = pptx.Presentation(f)
    f.close()
    slide = prs.slides[0]

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        p = text_frame.paragraphs[0]
        p.font.name = "Arial"
        p.font.bold = True
        p.font.size = Pt(115)
        p.text = name

    prs.save('updated.pptx')    

class MyHandler(BaseHTTPRequestHandler):
    def do_GET(self):
        print("Just received a GET request")

        qs = {}
        path = self.path
        if '?' in path:
            path, tmp = path.split('?', 1)
            qs = urlparse.parse_qs(tmp)

        try:
            ip = qs['ip'][0]
            essid = qs['essid'][0]
            print ip
            print essid
            self.send_response(200)
            f = open('ipfile.txt','w')
            f.write(ip + " " + essid)
            f.close()
        except Exception,e:
            self.send_response(200)
            self.wfile.write(str(e))
        return

        # f = None
        # try:
        #     name = qs['name'][0]
        #     name = name.replace("-","\n")
        #     generate_ppt(name)
        #     f = open('updated.pptx','rb')
        # except Exception,e:
        #     print str(e)


        # if f==None:
        #     self.wfile.write('Hello world')
        # else:
        #     self.send_response(200)
        #     self.send_header('Content-type','application/vnd.openxmlformats-officedocument.presentationml.presentation')
        #     self.send_header('Content-Disposition','attachment; filename='+name+'.pptx')

        #     self.end_headers()
        #     self.wfile.write(f.read())
        #     f.close()
        # return

    # def log_request(self, code=None, size=None):
    #     print('Request')

    # def log_message(self, format, *args):
    #     print('Message')

if __name__ == "__main__":
    try:
        server = HTTPServer(('', 8002), MyHandler)
        print('Started http server')
        server.serve_forever()
    except KeyboardInterrupt:
        print('^C received, shutting down server')
        server.socket.close()
